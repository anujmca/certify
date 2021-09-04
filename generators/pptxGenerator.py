import random
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
import os
from pptx import Presentation

def generate(template_file_path, data_sheet_dictionary):

    # data = {'employee_name': 'Anuj Kumar', 'position': '1'}

    ppt = Presentation(template_file_path)
    slide = ppt.slides[0]
    for shape in ppt.slides[0].shapes:
        if hasattr(shape, 'text') and shape.text[1:-1].lower() in data_sheet_dictionary:
            # shape.text = data_sheet_dictionary[shape.text.lower()]

            shape.text_frame.paragraphs[0].runs[0].text = data_sheet_dictionary[shape.text[1:-1].lower()]

            if len(shape.text_frame.paragraphs[0].runs) > 1:
                for index in range(1, len(shape.text_frame.paragraphs[0].runs), 1):
                    shape.text_frame.paragraphs[0].runs[index].text = ''

            # ppt.save(data_sheet_dictionary['employee_name'] + ".pptx")
    return ppt


def generate_in_memory(template_file, data_sheet_dictionary):

    # data = {'employee_name': 'Anuj Kumar', 'position': '1'}

    ppt = Presentation(template_file)
    slide = ppt.slides[0]
    for shape in ppt.slides[0].shapes:
        if hasattr(shape, 'text') and shape.text.lower() in data_sheet_dictionary:
            orig_text = shape.text

            # shape.text = data_sheet_dictionary[shape.text.lower()]
            shape.text = orig_text.replace(shape.text, str(data_sheet_dictionary[shape.text.lower()]))

    # ppt.save(data_sheet_dictionary['employee_name'] + ".pptx")
    return ppt