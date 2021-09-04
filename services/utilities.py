import os
from datetime import date

from pandas import *
from pptx import Presentation
import re

from enum import Enum
import uuid


# from django.db.models import Q
# # from pandas import *
# from services.models import *
# from services.models import Profile
from certify import settings


class BaseToken:
    id = 'id'
    first_name = 'first_name'
    last_name = 'last_name'
    phone_number = 'phone_number'
    email_id = 'email_id'


DEFAULT_TOKENS = [BaseToken.id, BaseToken.first_name, BaseToken.last_name, BaseToken.phone_number, BaseToken.email_id]


class Groups:
    issuer = 'issuer'
    awardee = 'awardee'


def get_user_group_names(user):
    return [group.name for group in user.groups.all()] if user is not None and user.groups is not None else []


def get_excel_headers(data_sheet):
    data_sheet_excel = ExcelFile(data_sheet)
    df = data_sheet_excel.parse(data_sheet_excel.sheet_names[0])
    df = df.reindex(sorted(df.columns), axis=1)
    return df.columns


def get_ppt_tokens(template_file):
    tokens = list()
    ppt = Presentation(template_file)
    slide = ppt.slides[0]
    for shape in ppt.slides[0].shapes:
        if hasattr(shape, 'text') and shape.text.startswith('{') and shape.text.endswith('}'):
            tokens.append(shape.text)

    tokens.sort()
    return tokens


def is_valid_email(email):
    regex_email = '^[a-z0-9]+[\._]?[a-z0-9]+[@]\w+[.]\w{2,3}$'
    try:
        if re.search(regex_email, email):
            return True
        else:
            return False
    except:
        return False


# import sys
# import argparse
# import os
# import fnmatch
# import win32com
#
# # import tkinter as Tkinter

from win32com.client import Dispatch
import pythoncom
from django.db import models
from django.core.files import File
import shutil
from django.core.files.storage import default_storage as storage


def get_jpg_file(file):
    ppt_dispatch = None
    prs = None
    thumbnail_file_path = None
    temp_ppt_file_path = None
    try:
        read_only = True
        has_title = False
        window = False
        unique_seed = str(uuid.uuid1())

        temp_dir = os.path.join(settings.MEDIA_ROOT, 'temp')
        os.makedirs(temp_dir, exist_ok=True)
        temp_ppt_file_path = os.path.join(temp_dir, unique_seed + file.name.split('/')[-1])

        shutil.copyfile(file.path, temp_ppt_file_path)

        pythoncom.CoInitialize()
        ppt_dispatch = Dispatch('Powerpoint.Application')
        prs = ppt_dispatch.Presentations.Open(temp_ppt_file_path, read_only, has_title, WithWindow=0)
        firstSlideRange = ppt_dispatch.Presentations[0].Slides.Range([1])

        thumbnails_dir = os.path.join(settings.MEDIA_ROOT, 'templates', 'thumbnails')
        os.makedirs(thumbnails_dir, exist_ok=True)
        thumbnail_file_path = os.path.join(thumbnails_dir, file.name.split('/')[-1].split('.')[0] + "_" + unique_seed + ".jpg")
        firstSlideRange.Export(thumbnail_file_path, 'JPG')
    except Exception as ex:
        thumbnail_file_path = None
    finally:
        if prs:
            prs.Close()
        if ppt_dispatch:
            ppt_dispatch.Quit()
        if temp_ppt_file_path and os.path.exists(temp_ppt_file_path):
            os.remove(temp_ppt_file_path)

    return thumbnail_file_path


def convert_ppt_to_pdf(ppt, pdf_filename):
    unique_seed = str(uuid.uuid1())
    temp_dir = os.path.join(settings.MEDIA_ROOT, 'temp')
    os.makedirs(temp_dir, exist_ok=True)
    temp_ppt_file_path = os.path.join(temp_dir, unique_seed + ".pptx")
    ppt.save(temp_ppt_file_path)

    todays_date = date.today()
    certificate_dir = os.path.join(settings.MEDIA_ROOT, 'certificates', str(todays_date.year), str(todays_date.month), str(todays_date.day))
    os.makedirs(certificate_dir, exist_ok=True)
    pdf_file_path = os.path.join(certificate_dir, pdf_filename)

    # return_pdf_file_path = PPTtoPDF(temp_ppt_file_path, pdf_file_path)

    ppt_dispatch = None
    deck = None
    try:
        pythoncom.CoInitialize()
        ppt_dispatch = Dispatch('Powerpoint.Application')
        deck = ppt_dispatch.Presentations.Open(temp_ppt_file_path, WithWindow=0)
        # ppt_dispatch.Presentations[0].SaveAs(pdf_file_path, 32)
        deck.SaveAs(pdf_file_path, 32)
    except:
        pdf_file_path = None
    finally:
        if deck:
            deck.Close()
        if ppt_dispatch:
            ppt_dispatch.Quit()
        os.remove(temp_ppt_file_path)
    return pdf_file_path



# import comtypes.client

# def PPTtoPDF(inputFileName, outputFileName, formatType = 32):
#     powerpoint = None
#     deck = None
#     return_pdf_file_path = None
#     try:
#         pythoncom.CoInitialize()
#         powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
#         # powerpoint.Visible = 0
#
#         if outputFileName[-3:] != 'pdf':
#             outputFileName = outputFileName + ".pdf"
#         deck = powerpoint.Presentations.Open(inputFileName)
#         deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
#         return_pdf_file_path =  outputFileName
#     except Exception as ex:
#         return_pdf_file_path = None
#     finally:
#         if deck:
#             deck.Close()
#         if powerpoint:
#             powerpoint.Quit()
#
#     return return_pdf_file_path